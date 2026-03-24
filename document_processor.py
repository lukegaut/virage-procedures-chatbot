"""
Processes .docx, .pdf, and .pptx procedure files into structured searchable data.
Extracts text organized by sections and saves embedded images.
Tables are processed in document order so they attach to the correct section.
"""

import os
import json
import hashlib
import re
from pathlib import Path
from docx import Document
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph
from lxml import etree

PROCEDURES_DIR = Path(__file__).parent / "procedures"
EXTRACTED_DIR = Path(__file__).parent / "extracted"
IMAGES_DIR = EXTRACTED_DIR / "images"
INDEX_FILE = EXTRACTED_DIR / "index.json"


def extract_images(doc, doc_id):
    """Extract all images from a document and save them, returning a mapping of rId -> filename."""
    image_map = {}
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_data = rel.target_part.blob
            ext = os.path.splitext(rel.target_part.partname)[1] or ".png"
            image_hash = hashlib.md5(image_data).hexdigest()[:8]
            filename = f"{doc_id}_{image_hash}{ext}"
            filepath = IMAGES_DIR / filename
            if not filepath.exists():
                with open(filepath, "wb") as f:
                    f.write(image_data)
            image_map[rel.rId] = filename
    return image_map


def get_paragraph_images(paragraph, image_map):
    """Find any images referenced in a paragraph."""
    images = []
    for run in paragraph.runs:
        xml = run._element.xml
        if "blip" in xml:
            rids = re.findall(r'r:embed="([^"]+)"', xml)
            for rid in rids:
                if rid in image_map:
                    images.append(image_map[rid])
    return images


def iter_block_items(doc):
    """
    Iterate over paragraphs AND tables in document order.
    This is critical — python-docx's doc.paragraphs and doc.tables
    are separate lists that lose ordering. This walks the XML tree
    to yield items in the order they appear in the document.
    """
    body = doc.element.body
    for child in body:
        tag = etree.QName(child).localname
        if tag == "p":
            yield Paragraph(child, doc)
        elif tag == "tbl":
            yield DocxTable(child, doc)


def get_table_images(table, image_map):
    """Extract images embedded inside table cells, with per-cell text context."""
    images = []
    image_cell_context = {}  # img_file -> cell text for context
    seen = set()
    for row in table.rows:
        row_text = " ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                para_imgs = get_paragraph_images(paragraph, image_map)
                for img in para_imgs:
                    if img not in seen:
                        seen.add(img)
                        images.append(img)
                        # Use the row text as context for this image
                        image_cell_context[img] = row_text
    return images, image_cell_context


def format_table(table):
    """Convert a docx table to readable text."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        # Deduplicate merged cells (python-docx repeats merged cell text)
        deduped = []
        prev = None
        for c in cells:
            if c != prev:
                deduped.append(c)
            prev = c
        rows.append(" | ".join(deduped))
    return "\n".join(rows)


def _add_after_context(section):
    """Add 'after' text context to images by scanning the content list."""
    if not section.get("image_contexts"):
        return
    content = section.get("content", [])
    # For each image, find text that comes after it in the content flow
    # Since images are inline with text, use the content lines after the "before" text
    for img_file, ctx in section["image_contexts"].items():
        before_text = ctx.get("before", "")
        # Find where in content the "before" text ends, then grab what follows
        found_idx = -1
        for i, line in enumerate(content):
            if before_text and line in before_text:
                found_idx = i
        # Grab up to 2 lines after
        after_lines = []
        if found_idx >= 0 and found_idx + 1 < len(content):
            after_lines = content[found_idx + 1: found_idx + 3]
        ctx["after"] = " ".join(after_lines) if after_lines else ""
        # Build a combined description for matching
        parts = []
        if ctx.get("parent"):
            parts.append(ctx["parent"])
        if ctx.get("section"):
            parts.append(ctx["section"])
        if ctx.get("before"):
            parts.append(ctx["before"])
        if ctx.get("after"):
            parts.append(ctx["after"])
        ctx["description"] = " ".join(parts)


def process_document(filepath):
    """Process a single .docx file into structured sections."""
    doc = Document(filepath)
    doc_name = Path(filepath).stem
    doc_id = hashlib.md5(doc_name.encode()).hexdigest()[:8]

    # Extract images
    image_map = extract_images(doc, doc_id)

    sections = []
    # Track parent headings so sub-sections inherit parent context
    parent_headings = {}  # level -> title

    current_section = {
        "title": doc_name,
        "level": 0,
        "content": [],
        "images": [],
        "image_contexts": {},  # filename -> surrounding text description
        "parent": "",
    }

    # We need to track recent text lines to build image context
    recent_text_lines = []

    for item in iter_block_items(doc):
        if isinstance(item, Paragraph):
            text = item.text.strip()
            if not text and not get_paragraph_images(item, image_map):
                continue

            style_name = item.style.name if item.style else ""
            para_images = get_paragraph_images(item, image_map)

            # Check if this is a heading
            is_heading = "Heading" in style_name or "heading" in style_name
            heading_level = 0
            if is_heading:
                try:
                    heading_level = int("".join(filter(str.isdigit, style_name)) or "1")
                except ValueError:
                    heading_level = 1

            if is_heading and text:
                # Save the current section if it has content
                if current_section["content"] or current_section["images"]:
                    sections.append(current_section)

                # Update parent heading tracking
                parent_headings[heading_level] = text
                # Clear any deeper level parents
                for lvl in list(parent_headings.keys()):
                    if lvl > heading_level:
                        del parent_headings[lvl]

                # Find parent title (closest heading with a lower level)
                parent = ""
                for lvl in sorted(parent_headings.keys()):
                    if lvl < heading_level:
                        parent = parent_headings[lvl]

                current_section = {
                    "title": text,
                    "level": heading_level,
                    "content": [],
                    "images": [],
                    "image_contexts": {},
                    "parent": parent,
                }
                recent_text_lines = []
            else:
                if text:
                    current_section["content"].append(text)
                    recent_text_lines.append(text)
                    # Keep only last 3 lines for context
                    if len(recent_text_lines) > 3:
                        recent_text_lines = recent_text_lines[-3:]
                if para_images:
                    current_section["images"].extend(para_images)
                    # Build context for each image from surrounding text
                    context_before = " ".join(recent_text_lines[-3:])
                    for img_file in para_images:
                        current_section["image_contexts"][img_file] = {
                            "before": context_before,
                            "section": current_section["title"],
                            "parent": current_section.get("parent", ""),
                        }

        elif isinstance(item, DocxTable):
            # Tables are now processed in order, attached to current section
            table_text = format_table(item)
            if table_text.strip():
                current_section["content"].append(f"[Table]\n{table_text}")
                recent_text_lines.append(table_text[:200])

            # Extract images embedded inside table cells
            table_imgs, cell_contexts = get_table_images(item, image_map)
            if table_imgs:
                current_section["images"].extend(table_imgs)
                for img_file in table_imgs:
                    current_section["image_contexts"][img_file] = {
                        "before": cell_contexts.get(img_file, ""),
                        "section": current_section["title"],
                        "parent": current_section.get("parent", ""),
                    }

    # Capture text AFTER images — go back and add "after" context
    for section in sections:
        _add_after_context(section)
    _add_after_context(current_section)

    # Don't forget the last section
    if current_section["content"] or current_section["images"]:
        sections.append(current_section)

    return {
        "filename": Path(filepath).name,
        "doc_name": doc_name,
        "doc_id": doc_id,
        "sections": sections,
    }


def _get_word_counts(text_lines, title):
    """Count meaningful words in a section's content, excluding stopwords and title words."""
    stopwords = {
        "the", "and", "for", "are", "but", "not", "you", "all", "can", "her",
        "was", "one", "our", "out", "its", "has", "his", "how", "man", "new",
        "now", "old", "see", "way", "who", "did", "get", "let", "say", "she",
        "too", "use", "when", "with", "from", "that", "this", "they", "will",
        "step", "crew", "2026", "pit", "stop", "mec", "puts", "goes", "takes",
        "makes", "sure", "position", "line", "behind", "white", "cross",
        "document", "their", "finish", "tasks", "orders", "them", "again",
        "stay", "when", "then", "into", "also", "before", "after",
    }
    title_words = set(w.lower() for w in re.split(r'\W+', title) if w)
    ignore = stopwords | title_words

    word_counts = {}
    for line in text_lines:
        for word in re.split(r'\W+', line.lower()):
            if len(word) >= 3 and word not in ignore and not word.isdigit():
                word_counts[word] = word_counts.get(word, 0) + 1
    return word_counts


def _disambiguate_sections(sections_group):
    """
    Given a list of sections with the same title, give each a unique suffix
    based on what makes it different from the others.
    """
    base_title = sections_group[0]["title"]

    # Get word counts for each section
    all_counts = []
    for section in sections_group:
        counts = _get_word_counts(section["content"], base_title)
        all_counts.append(counts)

    # For each section, find words that are distinctive to IT
    # (appear much more in this section than others)
    labels = {
        "gun": "Tyre Changes",
        "wheel": "Tyre Changes",
        "tyre": "Tyre Changes",
        "tire": "Tyre Changes",
        "tyres": "Tyre Changes",
        "electric": "Tyre Changes",
        "lance": "Tyre Changes",
        "driver": "Driver Change",
        "windshield": "Driver Change",
        "fuel": "Fuelling",
        "refuel": "Refuelling",
    }

    for i, section in enumerate(sections_group):
        my_counts = all_counts[i]
        other_counts = {}
        for j, counts in enumerate(all_counts):
            if j != i:
                for w, c in counts.items():
                    other_counts[w] = other_counts.get(w, 0) + c

        # Find words unique or dominant in this section
        distinctive = {}
        for word, count in my_counts.items():
            other_count = other_counts.get(word, 0)
            # Word appears significantly more in this section
            if count > other_count * 2 or (count >= 2 and other_count == 0):
                distinctive[word] = count

        # Try to match a known label
        best_label = None
        best_score = 0
        for word, count in distinctive.items():
            if word in labels and count > best_score:
                best_label = labels[word]
                best_score = count

        if best_label:
            section["title"] = f"{base_title} - {best_label}"
        else:
            # Fallback: use top distinctive words
            ranked = sorted(distinctive.items(), key=lambda x: -x[1])
            keywords = [w for w, c in ranked[:3] if c >= 2]
            if keywords:
                suffix = " ".join(w.capitalize() for w in keywords)
                section["title"] = f"{base_title} - {suffix}"
            else:
                section["title"] = f"{base_title} (Part {i + 1})"


def process_pdf(filepath):
    """
    Process a PDF file: render pages as images, grouped by major headings.
    Detects large/bold text as section headings and groups subsequent pages
    under them. When duplicate headings appear, disambiguates them using
    content keywords (e.g. "PIT STOP CREW LES 2026 - Tyre Changes").
    """
    import fitz  # PyMuPDF
    from PIL import Image
    import io

    doc_id = hashlib.md5(str(filepath).encode()).hexdigest()[:8]
    doc_name = Path(filepath).stem

    pdf = fitz.open(filepath)

    # First pass: extract page info and detect headings
    pages_info = []
    for page_num in range(len(pdf)):
        page = pdf[page_num]
        blocks = page.get_text("dict")["blocks"]

        text_lines = []
        heading = None
        max_font_size = 0

        for block in blocks:
            if block["type"] == 0:
                for line in block["lines"]:
                    line_text = "".join(span["text"] for span in line["spans"]).strip()
                    if not line_text or len(line_text) <= 1:
                        continue
                    font_size = max(span["size"] for span in line["spans"])
                    is_bold = any("bold" in span.get("font", "").lower() for span in line["spans"])

                    text_lines.append(line_text)

                    # Detect major headings (large bold text, typically 20px+)
                    if is_bold and font_size >= 20 and len(line_text) > 3:
                        if not heading or font_size > max_font_size:
                            heading = line_text
                            max_font_size = font_size

        # Render page as JPEG
        mat = fitz.Matrix(1.5, 1.5)
        pix = page.get_pixmap(matrix=mat)
        img_filename = f"{doc_id}_page{page_num + 1}.jpg"
        img_path = IMAGES_DIR / img_filename
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        img = img.convert("RGB")
        img.save(str(img_path), format="JPEG", quality=85)

        pages_info.append({
            "page_num": page_num + 1,
            "heading": heading,
            "text_lines": text_lines,
            "image": img_filename,
        })

    pdf.close()

    # Second pass: group pages by major headings into sections.
    # A "major heading" is one that differs from the current section's base title.
    # Sub-headings like "Step 1", "Step 2" stay within the current section.
    # Repeated major headings start a NEW section (disambiguated later).
    sections = []
    current_section = None
    current_base_title = None  # The major heading for the current section

    for page in pages_info:
        heading = page["heading"]

        # Decide if this heading starts a new section
        is_new_major_heading = False
        if heading:
            # "Step N" headings are sub-headings, not new sections
            is_step = bool(re.match(r'^Step\s+\d+', heading, re.IGNORECASE))
            if not is_step:
                if heading != current_base_title:
                    # Different heading than current — always new section
                    is_new_major_heading = True
                elif not page["text_lines"] or len(page["text_lines"]) <= 1:
                    # Same heading on a page with minimal text = title page for new section
                    is_new_major_heading = True

        if is_new_major_heading:
            if current_section:
                sections.append(current_section)
            current_base_title = heading
            current_section = {
                "title": heading,  # Will be disambiguated later if needed
                "base_title": heading,
                "level": 1,
                "content": [f"Document: {doc_name}"],
                "images": [page["image"]],
                "image_contexts": {},
                "parent": doc_name,
                "is_page_render": True,
            }
            current_section["content"].extend(page["text_lines"])
        elif current_section:
            # Append to current section
            current_section["images"].append(page["image"])
            current_section["content"].extend(page["text_lines"])
        else:
            # No section started yet — create a default one
            current_section = {
                "title": f"{doc_name} - Overview",
                "base_title": None,
                "level": 1,
                "content": [f"Document: {doc_name}"] + page["text_lines"],
                "images": [page["image"]],
                "image_contexts": {},
                "parent": doc_name,
                "is_page_render": True,
            }

    if current_section:
        sections.append(current_section)

    # Third pass: disambiguate sections with duplicate titles
    title_groups = {}
    for section in sections:
        t = section["title"]
        if t not in title_groups:
            title_groups[t] = []
        title_groups[t].append(section)

    for title, group in title_groups.items():
        if len(group) > 1:
            _disambiguate_sections(group)

    # Clean up: remove internal base_title field
    for section in sections:
        section.pop("base_title", None)

    return {
        "filename": Path(filepath).name,
        "doc_name": doc_name,
        "doc_id": doc_id,
        "sections": sections,
    }


def process_pptx(filepath):
    """Process a PowerPoint file into structured sections with images."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    doc_id = hashlib.md5(str(filepath).encode()).hexdigest()[:8]
    doc_name = Path(filepath).stem

    prs = Presentation(filepath)
    sections = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = ""
        slide_content = []
        slide_images = []
        image_contexts = {}

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # First text on slide with title placeholder is the title
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and not slide_title:
                            try:
                                if shape.placeholder_format.idx == 0:  # Title placeholder
                                    slide_title = text
                                    continue
                            except Exception:
                                pass
                        slide_content.append(text)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_data = shape.image.blob
                    ext = shape.image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    img_filename = f"{doc_id}_{hashlib.md5(img_data).hexdigest()[:8]}.{ext}"
                    img_path = IMAGES_DIR / img_filename
                    with open(img_path, "wb") as f:
                        f.write(img_data)
                    slide_images.append(img_filename)
                    context_before = " ".join(slide_content[-3:]) if slide_content else slide_title
                    image_contexts[img_filename] = {
                        "before": context_before,
                        "section": slide_title or f"Slide {slide_num}",
                        "parent": doc_name,
                    }
                except Exception:
                    pass

            # Handle tables in slides
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    deduped = []
                    prev = None
                    for c in cells:
                        if c != prev:
                            deduped.append(c)
                        prev = c
                    rows.append(" | ".join(deduped))
                slide_content.append(f"[Table]\n" + "\n".join(rows))

        if slide_content or slide_images:
            sections.append({
                "title": slide_title or f"Slide {slide_num}",
                "level": 1,
                "content": slide_content,
                "images": slide_images,
                "image_contexts": image_contexts,
                "parent": doc_name,
            })

    return {
        "filename": Path(filepath).name,
        "doc_name": doc_name,
        "doc_id": doc_id,
        "sections": sections,
    }


def build_index():
    """Process all procedure files (.docx, .pdf, .pptx) in the procedures directory and build the search index."""
    EXTRACTED_DIR.mkdir(parents=True, exist_ok=True)
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)

    index = {"documents": []}

    # Find all supported file types
    all_files = []
    for ext in ("*.docx", "*.pdf", "*.pptx"):
        all_files.extend(PROCEDURES_DIR.glob(ext))

    if not all_files:
        print(f"No procedure files found in {PROCEDURES_DIR}")
        print(f"Supported formats: .docx, .pdf, .pptx")
        return index

    for filepath in all_files:
        print(f"Processing: {filepath.name}")
        try:
            ext = filepath.suffix.lower()
            if ext == ".docx":
                doc_data = process_document(filepath)
            elif ext == ".pdf":
                doc_data = process_pdf(filepath)
            elif ext == ".pptx":
                doc_data = process_pptx(filepath)
            else:
                continue
            index["documents"].append(doc_data)
            print(f"  -> {len(doc_data['sections'])} sections extracted")
        except Exception as e:
            print(f"  -> Error processing {filepath.name}: {e}")

    # Save index
    with open(INDEX_FILE, "w", encoding="utf-8") as f:
        json.dump(index, f, indent=2, ensure_ascii=False)

    print(f"\nIndex saved to {INDEX_FILE}")
    print(f"Images saved to {IMAGES_DIR}")
    return index


if __name__ == "__main__":
    build_index()
